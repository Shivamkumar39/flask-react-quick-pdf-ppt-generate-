from flask import Flask, request, send_file, jsonify
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Pt
import io
import os
import re

app = Flask(__name__)
# Assuming CORS enabled as before

def safe_filename(name: str) -> str:
    return re.sub(r'[^a-zA-Z0-9_\-]', '_', name)

def create_ppt_bytes(name, topic, description, slide_count):
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = topic or "Presentation"
    subtitle = f"Presented by {name}" if name else ""
    try:
        slide.placeholders[1].text = subtitle
    except Exception:
        pass

    content_slide_layout = prs.slide_layouts[1]
    points = []
    if description:
        if '.' in description:
            points = [p.strip() for p in description.split('.') if p.strip()]
        elif ',' in description:
            points = [p.strip() for p in description.split(',') if p.strip()]
        else:
            points = [description]

    for i in range(max(slide_count, 1)):
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = f"Slide {i+1}"
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        bullet_points = points[i * 5 : (i + 1) * 5] if points else [f"Point {i+1}"]

        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

@app.route('/generate-pdf', methods=['POST'])
def generate_pdf():
    data = request.json
    name = data.get('name')
    topic = data.get('topic')
    description = data.get('description')
    slide_count = int(data.get('slide_count', 1))

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt=f"Topic: {topic}", ln=True, align='C')
    pdf.cell(200, 10, txt=f"Presented by: {name}", ln=True, align='C')
    pdf.multi_cell(0, 10, txt=f"Description:\n{description}")

    for i in range(slide_count):
        pdf.add_page()
        pdf.set_font("Arial", size=16)
        pdf.cell(200, 10, txt=f"Slide {i+1}", ln=True, align='C')

    output = io.BytesIO()
    pdf.output(output)
    output.seek(0)
    safe_topic = safe_filename(topic) or "presentation"
    return send_file(output, download_name=f"{safe_topic}.pdf", as_attachment=True, mimetype="application/pdf")


@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    data = request.json
    name = data.get('name')
    topic = data.get('topic')
    description = data.get('description')
    slide_count = int(data.get('slide_count', 1))

    ppt_io = create_ppt_bytes(name, topic, description, slide_count)
    safe_topic = safe_filename(topic) or "presentation"
    return send_file(ppt_io, download_name=f"{safe_topic}.pptx", as_attachment=True, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

if __name__ == "__main__":
    app.run(debug=True)
