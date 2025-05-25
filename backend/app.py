# import os
# import re
# import traceback
# from flask import Flask, request, jsonify
# from flask_cors import CORS
# from pptx import Presentation
# from fpdf import FPDF
# from dotenv import load_dotenv
# from google import generativeai as genai
# from fpdf import FPDF

# load_dotenv()

# API_KEY = os.getenv("GOOGLE_API_KEY")

# app = Flask(__name__)
# CORS(app)
# UPLOAD_FOLDER = "static"
# os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# genai.configure(api_key=API_KEY)

# # --- CHANGE THIS LINE ---
# # Try 'gemini-1.5-flash' first. If it still fails, try 'gemini-pro'.
# MODEL = "gemini-1.5-flash" 
# # MODEL = "gemini-pro" # Uncomment this line and comment the above if 'gemini-1.5-flash' fails too
# def safe_filename(name: str) -> str:
#     """Sanitize the filename by replacing unsafe characters with underscores."""
#     return re.sub(r'[^a-zA-Z0-9_\-]', '_', name)

# def generate_text_with_gemini(topic: str, description: str, slides: int) -> str:
#     """
#     Generate a professional slide-based presentation using Gemini.

#     Args:
#         topic (str): The topic of the presentation.
#         description (str): A short description or context for the presentation.
#         slides (int): Number of slides to generate.

#     Returns:
#         str: Formatted text containing slides with titles and bullet points.
#     """
#     prompt = (
#         f"Create a professional presentation consisting of {slides} slides.\n"
#         f"Topic: {topic}\n"
#         f"Description: {description}\n\n"
#         "Each slide must include a title (first line) and bullet points (each starting with '-').\n"
#         "Separate each slide with two newlines.\n"
#         "Ensure the content is informative, concise, and suitable for professional use."
#     )

#     model = genai.GenerativeModel(MODEL)
#     response = model.generate_content(
#         prompt,
#         generation_config={
#             "temperature": 0.7,
#             "top_p": 0.8,
#             "max_output_tokens": 2048,
#         },
#         safety_settings=[
#             {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
#             {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
#             {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
#             {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
#         ]
#     )

#     return response.text
# # ... (rest of your generate_ppt, generate_pdf, and @app.route("/generate") functions remain the same) ...

# def generate_ppt(content: str, filename: str) -> str:
#     prs = Presentation()
#     slides = [s.strip() for s in content.split("\n\n") if s.strip()]

#     for slide_text in slides:
#         slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
#         lines = slide_text.split("\n")
#         title = lines[0] if lines else "Slide"
#         body_lines = [line.strip() for line in lines[1:] if line.strip()]

#         slide.shapes.title.text = title
#         text_frame = slide.placeholders[1].text_frame
#         text_frame.clear()

#         for i, line in enumerate(body_lines):
#             bullet = line[1:].strip() if line.startswith("-") else line
#             p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
#             p.text = bullet
#             p.level = 0

#     ppt_path = os.path.join(UPLOAD_FOLDER, filename)
#     prs.save(ppt_path)
#     return ppt_path

# UPLOAD_FOLDER = "pdfs"
# os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# FONT_PATH = "fonts/DejaVuSans.ttf"

# def generate_pdf(content: str, filename: str) -> str:
#     pdf = FPDF()
#     pdf.add_page()
#     pdf.set_auto_page_break(auto=True, margin=15)

#     # Use DejaVuSans for full Unicode support
#     pdf.add_font("DejaVu", "", FONT_PATH, uni=True)
#     pdf.set_font("DejaVu", size=12)

#     for para in content.split("\n\n"):
#         para = para.strip()
#         if para:
#             pdf.multi_cell(0, 10, para)
#             pdf.ln(2)

#     pdf_path = os.path.join(UPLOAD_FOLDER, filename)
#     pdf.output(pdf_path)
#     return pdf_path
# @app.route("/generate", methods=["POST"])
# def generate_doc():
#     data = request.json or {}

#     name = data.get("name", "").strip()
#     to = data.get("to", "").strip()
#     from_ = data.get("from", "").strip()
#     title = data.get("title", "").strip()
#     message = data.get("message", "").strip()
#     date = data.get("date", "").strip()
#     doc_type = data.get("doc_type", "Letter")
#     style = data.get("style", "Classic")

#     if not all([name, to, from_, title, message, date]):
#         return jsonify({"error": "All fields are required."}), 400

#     try:
#         prompt = (
#             f"Generate a {style.lower()} styled {doc_type.lower()} with the following details:\n"
#             f"- To: {to}\n"
#             f"- From: {from_}\n"
#             f"- Title: {title}\n"
#             f"- Message: {message}\n"
#             f"- Date: {date}\n"
#             f"Ensure the output is professional and formatted appropriately for a {doc_type}.\n"
#             f"Use clear paragraphs and structure."
#         )

#         model = genai.GenerativeModel(MODEL)
#         response = model.generate_content(
#             prompt,
#             generation_config={"temperature": 0.7, "top_p": 0.9, "max_output_tokens": 2048}
#         )

#         content = response.text
#         safe_name = safe_filename(name)
#         pdf_file = f"{safe_name}_{doc_type.lower()}.pdf"
#         ppt_file = f"{safe_name}_{doc_type.lower()}.pptx"

#         generate_pdf(content, pdf_file)
#         generate_ppt(content, ppt_file)

#         return jsonify({
#             "pdf_url": f"/static/{pdf_file}",
#             "ppt_url": f"/static/{ppt_file}"
#         })

#     except Exception as e:
#         traceback.print_exc()
#         return jsonify({"error": str(e)}), 500



# if __name__ == "__main__":
#     app.run(debug=True)


import os
import re
import traceback
import requests
from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from fpdf import FPDF
from dotenv import load_dotenv
from google import generativeai as genai

load_dotenv()

API_KEY = os.getenv("GOOGLE_API_KEY")

app = Flask(__name__)
CORS(app)

# Use one folder for uploads/static files
UPLOAD_FOLDER = "static"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Font setup
FONT_DIR = "fonts"
os.makedirs(FONT_DIR, exist_ok=True)
FONT_PATH = os.path.join(FONT_DIR, "DejaVuSans.ttf")
FONT_URL = "https://github.com/dejavu-fonts/dejavu-fonts/raw/master/ttf/DejaVuSans.ttf"

def download_font():
    if not os.path.exists(FONT_PATH):
        print("Downloading DejaVuSans.ttf font...")
        response = requests.get(FONT_URL)
        if response.status_code == 200:
            with open(FONT_PATH, "wb") as f:
                f.write(response.content)
            print("Font downloaded successfully.")
        else:
            print(f"Failed to download font: HTTP {response.status_code}")
    else:
        print("Font already exists.")

download_font()

genai.configure(api_key=API_KEY)

# --- CHANGE THIS LINE ---
# Try 'gemini-1.5-flash' first. If it still fails, try 'gemini-pro'.
MODEL = "gemini-1.5-flash"
# MODEL = "gemini-pro" # Uncomment if above fails

def safe_filename(name: str) -> str:
    """Sanitize the filename by replacing unsafe characters with underscores."""
    return re.sub(r'[^a-zA-Z0-9_\-]', '_', name)

def generate_text_with_gemini(topic: str, description: str, slides: int) -> str:
    prompt = (
        f"Create a professional presentation consisting of {slides} slides.\n"
        f"Topic: {topic}\n"
        f"Description: {description}\n\n"
        "Each slide must include a title (first line) and bullet points (each starting with '-').\n"
        "Separate each slide with two newlines.\n"
        "Ensure the content is informative, concise, and suitable for professional use."
    )
    model = genai.GenerativeModel(MODEL)
    response = model.generate_content(
        prompt,
        generation_config={
            "temperature": 0.7,
            "top_p": 0.8,
            "max_output_tokens": 2048,
        },
        safety_settings=[
            {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
            {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
            {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
            {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
        ]
    )
    return response.text

def generate_ppt(content: str, filename: str) -> str:
    prs = Presentation()
    slides = [s.strip() for s in content.split("\n\n") if s.strip()]

    for slide_text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        lines = slide_text.split("\n")
        title = lines[0] if lines else "Slide"
        body_lines = [line.strip() for line in lines[1:] if line.strip()]

        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for i, line in enumerate(body_lines):
            bullet = line[1:].strip() if line.startswith("-") else line
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = bullet
            p.level = 0

    ppt_path = os.path.join(UPLOAD_FOLDER, filename)
    prs.save(ppt_path)
    return ppt_path

def generate_pdf(content: str, filename: str) -> str:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Use DejaVuSans for full Unicode support
    pdf.add_font("DejaVu", "", FONT_PATH, uni=True)
    pdf.set_font("DejaVu", size=12)

    for para in content.split("\n\n"):
        para = para.strip()
        if para:
            pdf.multi_cell(0, 10, para)
            pdf.ln(2)

    pdf_path = os.path.join(UPLOAD_FOLDER, filename)
    pdf.output(pdf_path)
    return pdf_path

@app.route("/generate", methods=["POST"])
def generate_doc():
    data = request.json or {}

    name = data.get("name", "").strip()
    to = data.get("to", "").strip()
    from_ = data.get("from", "").strip()
    title = data.get("title", "").strip()
    message = data.get("message", "").strip()
    date = data.get("date", "").strip()
    doc_type = data.get("doc_type", "Letter")
    style = data.get("style", "Classic")

    if not all([name, to, from_, title, message, date]):
        return jsonify({"error": "All fields are required."}), 400

    try:
        prompt = (
            f"Generate a {style.lower()} styled {doc_type.lower()} with the following details:\n"
            f"- To: {to}\n"
            f"- From: {from_}\n"
            f"- Title: {title}\n"
            f"- Message: {message}\n"
            f"- Date: {date}\n"
            f"Ensure the output is professional and formatted appropriately for a {doc_type}.\n"
            f"Use clear paragraphs and structure."
        )

        model = genai.GenerativeModel(MODEL)
        response = model.generate_content(
            prompt,
            generation_config={"temperature": 0.7, "top_p": 0.9, "max_output_tokens": 2048}
        )

        content = response.text
        safe_name = safe_filename(name)
        pdf_file = f"{safe_name}_{doc_type.lower()}.pdf"
        ppt_file = f"{safe_name}_{doc_type.lower()}.pptx"

        generate_pdf(content, pdf_file)
        generate_ppt(content, ppt_file)

        return jsonify({
            "pdf_url": f"/static/{pdf_file}",
            "ppt_url": f"/static/{ppt_file}"
        })

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    app.run(debug=True)
